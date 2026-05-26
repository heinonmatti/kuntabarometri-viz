"""Convert raw askia JSON + per-kunta PPTX into one tidy CSV.

Output schema (long format, one row per data point)::

    source_type        timeseries  | cross_section
    kind               theme       | sub_question
    indicator_id       e.g. "Kokonaisranking", "osio1", "subq_03"
    indicator_label_fi human-readable question / theme name
    theme_label_fi     theme this sub-question belongs to (NA for theme rows)
    year               int (2020/2022/2024/2026 for theme rows; 2026 for cross-section)
    kunta_slug         hamina | kotka | kouvola | kokomaa
    kunta_label_fi     "Hamina" | "Kotka" | "Kouvola" | "Koko maa"
    mean               weighted mean of Likert 1-5
    n                  count of respondents (sum of Likert counts excluding "En osaa sanoa")
    sd                 standard deviation of the Likert distribution
    se                 standard error of the mean (sd / sqrt(n))
    ci95               half-width of the 95% confidence interval (1.96 * se)
"""

from __future__ import annotations

import csv
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation

LIKERT_SCORES = (1, 2, 3, 4, 5)
EOS_PATTERN = re.compile(r"^en\s*osaa\s*sanoa", re.IGNORECASE)
TRAILING_YEAR_PATTERN = re.compile(r"\s*\(\d{4}\)\s*$")


@dataclass(frozen=True)
class Row:
    source_type: str
    kind: str
    indicator_id: str
    indicator_label_fi: str
    theme_label_fi: str | None
    year: int
    kunta_slug: str
    kunta_label_fi: str
    mean: float | None
    n: int
    sd: float | None = None
    se: float | None = None
    ci95: float | None = None
    scale: str = "likert_1_5"  # also: "nps" (range -100..+100), "other"


def _mean_sd_se_from_counts(counts: dict[int, float]) -> tuple[float | None, float | None, float | None, float | None, int]:
    """Return (mean, sd, se, ci95, n) for a {score: count} dict."""
    n_total = sum(counts.values())
    if n_total <= 0:
        return None, None, None, None, 0
    mean = sum(score * c for score, c in counts.items()) / n_total
    variance = sum(((score - mean) ** 2) * c for score, c in counts.items()) / n_total
    sd = variance ** 0.5
    # SE of the mean uses n - we round to int respondents
    n_int = int(round(n_total))
    se = sd / (n_int ** 0.5) if n_int > 0 else None
    ci95 = 1.96 * se if se is not None else None
    return mean, sd, se, ci95, n_int


# ---------------------------------------------------------------------------
# AskiaVista JSON -> theme-level time series rows
# ---------------------------------------------------------------------------


def _category_stats(chart: dict) -> list[tuple[str, dict[int, float]]]:
    """Return [(category_name, {score: count}), ...] for one askia chart payload.

    Lets the caller compute mean/SD/CI from the per-Likert distribution and
    also aggregate across categories (e.g. national Koko maa pool).
    """
    cats = [c.get("name") if isinstance(c, dict) else c for c in chart.get("categories", [])]
    series = chart.get("series", [])
    likert_series = [s for s in series if isinstance(s, dict) and _safe_score(s.get("name")) is not None]
    out: list[tuple[str, dict[int, float]]] = []
    for i, name in enumerate(cats):
        counts: dict[int, float] = {}
        for s in likert_series:
            score = _safe_score(s.get("name"))
            cnt = s["data"][i] if i < len(s.get("data", [])) else 0
            counts[score] = counts.get(score, 0) + cnt
        out.append((name or "", counts))
    return out


def _safe_score(name) -> int | None:
    """Extract the leading Likert score from a series label.

    Accepts 1-5 (standard Likert).
    """
    if name is None:
        return None
    m = re.match(r"^\s*([1-5])\b", str(name))
    if not m:
        return None
    return int(m.group(1))


def _safe_nps_score(name) -> int | None:
    """Extract a 0–10 NPS rating from a series label."""
    if name is None:
        return None
    m = re.match(r"^\s*(\d{1,2})\b", str(name))
    if not m:
        return None
    s = int(m.group(1))
    return s if 0 <= s <= 10 else None


def _kokomaa_stats(category_stats: list[tuple[str, dict[int, float]]]) -> tuple[float | None, float | None, float | None, float | None, int]:
    """Pool counts across all individual kunta rows (excludes 'Keskiarvo (...)' rows)."""
    pooled: dict[int, float] = {}
    for name, counts in category_stats:
        if not name or name.startswith("Keskiarvo"):
            continue
        for k, v in counts.items():
            pooled[k] = pooled.get(k, 0) + v
    return _mean_sd_se_from_counts(pooled)


def _find_kunta_stats(category_stats: list[tuple[str, dict[int, float]]], kunta_label: str) -> tuple[float | None, float | None, float | None, float | None, int]:
    for name, counts in category_stats:
        if name and (name == kunta_label or name.startswith(kunta_label + " ")):
            return _mean_sd_se_from_counts(counts)
    return None, None, None, None, 0


def askia_rows(
    json_path: Path,
    theme: dict,
    year: int,
    focal: dict,
    comparators: list[dict],
    respondent_counts: dict | None = None,
) -> Iterable[Row]:
    """Build theme-level time-series rows.

    The *mean* uses answer-level pooling (sum over the per-Likert answer
    counts in the askia chart). The *SE / CI95* however use the true
    **respondent count** for that (kunta, year), passed in via
    ``respondent_counts``. Without that correction CIs would be ~1.5–6×
    too narrow because every respondent contributes multiple correlated
    answers per theme.
    """
    payload = json.loads(json_path.read_text(encoding="utf-8"))
    chart_item = next((c for c in payload if isinstance(c, dict) and c.get("type") == "chart"), None)
    if chart_item is None or not chart_item.get("output"):
        return
    chart = chart_item["output"][0]
    cat_stats = _category_stats(chart)

    def _resp_n(slug: str) -> int | None:
        if not respondent_counts:
            return None
        per_year = respondent_counts.get(str(year), {})
        n = per_year.get(slug)
        return int(n) if n else None

    def _row(kunta_slug: str, kunta_label: str, stats: tuple) -> Row:
        mean, sd, se, ci95, n_answers = stats
        resp_n = _resp_n(kunta_slug)
        if resp_n and resp_n > 0 and sd is not None:
            se = sd / (resp_n ** 0.5)
            ci95 = 1.96 * se
            n_eff = resp_n
        else:
            n_eff = n_answers
        return Row(
            source_type="timeseries",
            kind="theme",
            indicator_id=theme["askia_rows"],
            indicator_label_fi=theme["label_fi"],
            theme_label_fi=theme["label_fi"],
            year=year,
            kunta_slug=kunta_slug,
            kunta_label_fi=kunta_label,
            mean=mean,
            n=n_eff,
            sd=sd,
            se=se,
            ci95=ci95,
        )

    yield _row(focal["slug"], focal["label"], _find_kunta_stats(cat_stats, focal["label"]))
    for comp in comparators:
        stats = _kokomaa_stats(cat_stats) if comp.get("is_aggregate") else _find_kunta_stats(cat_stats, comp["label"])
        yield _row(comp["slug"], comp["label"], stats)


# ---------------------------------------------------------------------------
# PPTX -> 2026 sub-question cross-section rows
# ---------------------------------------------------------------------------


def _strip_year_suffix(s: str) -> str:
    return TRAILING_YEAR_PATTERN.sub("", s).strip()


def _nps_stats_from_counts_0to10(counts: dict[int, float]) -> tuple[float, float, float, float, int]:
    """Net Promoter Score from a 0–10 rating distribution.

    Each respondent contributes +1 (promoter, 9–10), 0 (passive, 7–8) or
    −1 (detractor, 0–6). NPS = mean × 100, reported on the −100…+100 axis.
    SE is computed from the per-respondent {−1, 0, +1} distribution.
    """
    n_total = sum(counts.values())
    if n_total <= 0:
        return None, None, None, None, 0
    nps_values: dict[float, float] = {}
    for rating, cnt in counts.items():
        v = 1.0 if rating >= 9 else (-1.0 if rating <= 6 else 0.0)
        nps_values[v] = nps_values.get(v, 0) + cnt
    mean_per_r = sum(v * c for v, c in nps_values.items()) / n_total
    variance = sum(((v - mean_per_r) ** 2) * c for v, c in nps_values.items()) / n_total
    sd_per_r = variance ** 0.5
    n_int = int(round(n_total))
    se_per_r = sd_per_r / (n_int ** 0.5) if n_int > 0 else None
    # Scale to NPS units (×100)
    return (
        mean_per_r * 100,
        sd_per_r * 100,
        se_per_r * 100 if se_per_r is not None else None,
        1.96 * se_per_r * 100 if se_per_r is not None else None,
        n_int,
    )


def _pptx_chart_focal_stats(chart) -> tuple[float | None, float | None, float | None, float | None, int, str | None, str]:
    """Return (mean, sd, se, ci95, n, focal_cat_label, scale) for the first category.

    PPTX bar-stacked charts have one series per rating score. For 1–5 Likert
    questions the series are named "5 Kehittynyt ...", …, "1 Kehittynyt ...".
    For the single 0–10 NPS question the series are "10 Suosittelisin
    varmasti", "9", …, "0 En suosittelisi lainkaan". We detect the scale by
    the set of leading score numbers and compute mean/SD/SE/CI
    appropriately. The returned ``scale`` is one of "likert_1_5" or "nps".
    """
    try:
        plot = chart.plots[0]
        cats = [str(c) for c in plot.categories]
    except Exception:
        return None, None, None, None, 0, None, "likert_1_5"
    if not cats:
        return None, None, None, None, 0, None, "likert_1_5"

    focal_idx = 0
    cat = cats[focal_idx]
    n_match = re.search(r"n\s*=\s*([\d\s]+)", cat)
    n_total = int(n_match.group(1).replace(" ", "")) if n_match else 0
    if n_total == 0:
        return None, None, None, None, 0, cat, "likert_1_5"

    # Detect scale by series names
    score_set: set[int] = set()
    for s in chart.series:
        s_nps = _safe_nps_score(s.name)
        if s_nps is not None:
            score_set.add(s_nps)
    is_nps = max(score_set, default=-1) >= 6 and min(score_set, default=99) == 0
    scale = "nps" if is_nps else "likert_1_5"
    parser = _safe_nps_score if is_nps else _safe_score

    counts: dict[int, float] = {}
    for s in chart.series:
        score = parser(s.name)
        if score is None:
            continue
        try:
            vals = list(s.values)
        except Exception:
            continue
        if focal_idx >= len(vals) or vals[focal_idx] is None:
            continue
        pct = float(vals[focal_idx])
        counts[score] = counts.get(score, 0) + (pct / 100.0) * n_total
    if is_nps:
        mean, sd, se, ci95, n = _nps_stats_from_counts_0to10(counts)
    else:
        mean, sd, se, ci95, n = _mean_sd_se_from_counts(counts)
    if mean is None:
        m = re.search(r"ka=([-\d,\.]+)", cat)
        if m:
            mean = float(m.group(1).replace(",", "."))
            n = n_total
    return mean, sd, se, ci95, n, cat, scale


def pptx_rows(pptx_path: Path, kunta_slug: str, kunta_label: str, year: int = 2026) -> Iterable[Row]:
    prs = Presentation(pptx_path)
    seen_titles: dict[str, int] = {}
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not (hasattr(shape, "has_chart") and shape.has_chart):
                continue
            ch = shape.chart
            if not (hasattr(ch, "has_title") and ch.has_title):
                continue
            try:
                title = ch.chart_title.text_frame.text.strip()
            except Exception:
                title = ""
            if not title:
                continue
            label = _strip_year_suffix(title)
            if not label:
                continue
            mean, sd, se, ci95, n, _, scale = _pptx_chart_focal_stats(ch)
            if mean is None:
                continue
            base_id = re.sub(r"[^a-z0-9]+", "_", label.lower())[:60].strip("_")
            count = seen_titles.get(base_id, 0)
            seen_titles[base_id] = count + 1
            indicator_id = f"{base_id}_{count}" if count else base_id
            yield Row(
                source_type="cross_section",
                kind="sub_question",
                indicator_id=f"subq_{slide_idx:02d}_{indicator_id}",
                indicator_label_fi=label,
                theme_label_fi=None,
                year=year,
                kunta_slug=kunta_slug,
                kunta_label_fi=kunta_label,
                mean=mean,
                n=n,
                sd=sd,
                se=se,
                ci95=ci95,
                scale=scale,
            )


# ---------------------------------------------------------------------------
# Combine everything
# ---------------------------------------------------------------------------


def write_csv(rows: Iterable[Row], dest: Path) -> int:
    dest.parent.mkdir(parents=True, exist_ok=True)
    fieldnames = [
        "source_type",
        "kind",
        "indicator_id",
        "indicator_label_fi",
        "theme_label_fi",
        "year",
        "kunta_slug",
        "kunta_label_fi",
        "mean",
        "n",
        "sd",
        "se",
        "ci95",
        "scale",
    ]
    count = 0

    def _fmt(v):
        return "" if v is None else f"{v:.6f}"

    with dest.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for r in rows:
            writer.writerow(
                {
                    "source_type": r.source_type,
                    "kind": r.kind,
                    "indicator_id": r.indicator_id,
                    "indicator_label_fi": r.indicator_label_fi,
                    "theme_label_fi": r.theme_label_fi or "",
                    "year": r.year,
                    "kunta_slug": r.kunta_slug,
                    "kunta_label_fi": r.kunta_label_fi,
                    "mean": _fmt(r.mean),
                    "n": r.n,
                    "sd": _fmt(r.sd),
                    "se": _fmt(r.se),
                    "ci95": _fmt(r.ci95),
                    "scale": r.scale,
                }
            )
            count += 1
    return count


def build_all(cfg: dict, raw_dir: Path, processed_dir: Path) -> Path:
    focal = cfg["focal_kunta"]
    comparators = cfg["comparators"]
    themes = cfg["themes"]
    years = cfg["years"]

    # Respondent counts (n per kunta per year) — used as the SE denominator
    # on askia themes to avoid the answer-clustering bias.
    resp_path = raw_dir / "askia" / "_respondent_counts.json"
    respondent_counts = json.loads(resp_path.read_text(encoding="utf-8")) if resp_path.exists() else None

    # Theme labels covered by askia time-series — used to suppress PPTX
    # summary slides for the same theme (otherwise the dashboard shows
    # two cards with the same title and slightly different numbers).
    theme_labels = {t["label_fi"] for t in themes}

    rows: list[Row] = []

    # 1) askia theme-level time series (mean + respondent-n SE)
    askia_dir = raw_dir / "askia"
    for theme in themes:
        for year in years:
            jp = askia_dir / f"{theme['askia_rows']}_{year}.json"
            if not jp.exists():
                continue
            rows.extend(askia_rows(jp, theme, year, focal, comparators, respondent_counts))

    # 2) PPTX cross-section, skipping theme-summary slides that duplicate askia
    pptx_dir = raw_dir / "pptx"
    kuntas = [focal] + comparators
    for k in kuntas:
        path = pptx_dir / f"{k['slug']}.pptx"
        if not path.exists():
            continue
        for r in pptx_rows(path, k["slug"], k["label"]):
            if r.indicator_label_fi in theme_labels:
                continue
            rows.append(r)

    dest = processed_dir / "kuntabarometri.csv"
    n = write_csv(rows, dest)
    print(f"Wrote {n:,} rows to {dest}")
    return dest
